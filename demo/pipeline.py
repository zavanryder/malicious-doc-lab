"""Document processing pipeline: parse -> OCR -> chunk -> embed -> store."""

import io
from email import policy
from email.parser import BytesParser
from pathlib import Path

import chromadb
import fitz  # PyMuPDF
import ollama
from bs4 import BeautifulSoup
from docx import Document
from PIL import Image
import pytesseract

from config import (
    CHUNK_OVERLAP,
    CHUNK_SIZE,
    CHROMA_COLLECTION,
    OLLAMA_BASE_URL,
    OLLAMA_EMBED_MODEL,
)

# Global state
chroma_client = chromadb.Client()
collection = chroma_client.get_or_create_collection(name=CHROMA_COLLECTION)
ollama_client = ollama.Client(host=OLLAMA_BASE_URL)

# Store last-processed document state for evidence endpoints
last_extracted_text: str = ""
last_chunks: list[str] = []


def reset_state():
    """Clear all ingested data."""
    global last_extracted_text, last_chunks, collection
    last_extracted_text = ""
    last_chunks = []
    chroma_client.delete_collection(name=CHROMA_COLLECTION)
    collection = chroma_client.get_or_create_collection(name=CHROMA_COLLECTION)


def parse_document(file_bytes: bytes, filename: str) -> str:
    """Extract text from a document based on file extension."""
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(file_bytes)
    elif suffix == ".docx":
        return _parse_docx(file_bytes)
    elif suffix == ".xlsx":
        return _parse_xlsx(file_bytes)
    elif suffix == ".pptx":
        return _parse_pptx(file_bytes)
    elif suffix == ".eml":
        return _parse_eml(file_bytes)
    elif suffix in (".html", ".htm"):
        return _parse_html(file_bytes)
    elif suffix in (".csv", ".txt", ".md"):
        return file_bytes.decode("utf-8", errors="replace")
    elif suffix in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
        return _ocr_image(file_bytes)
    else:
        return file_bytes.decode("utf-8", errors="replace")


def _parse_pdf(data: bytes) -> str:
    """Extract text from PDF, including OCR on embedded images."""
    doc = fitz.open(stream=data, filetype="pdf")
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
        # Extract and OCR any images on the page
        for img_info in page.get_images():
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            if base_image:
                img = Image.open(io.BytesIO(base_image["image"]))
                ocr_text = pytesseract.image_to_string(img)
                if ocr_text.strip():
                    text_parts.append(ocr_text)
    return "\n".join(text_parts)


def _parse_docx(data: bytes) -> str:
    """Extract text from DOCX, including metadata."""
    doc = Document(io.BytesIO(data))
    parts = []
    # Extract metadata (intentionally vulnerable — includes all fields)
    props = doc.core_properties
    for attr in ("author", "subject", "keywords", "description", "title"):
        val = getattr(props, attr, None)
        if val:
            parts.append(f"{attr}: {val}")
    for para in doc.paragraphs:
        parts.append(para.text)
    return "\n".join(parts)


def _parse_xlsx(data: bytes) -> str:
    """Extract text and metadata from XLSX workbooks."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []

    props = wb.properties
    for attr in ("creator", "subject", "keywords", "description", "title"):
        val = getattr(props, attr, None)
        if val:
            parts.append(f"{attr}: {val}")

    for sheet in wb.worksheets:
        parts.append(f"[sheet] {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def _parse_pptx(data: bytes) -> str:
    """Extract text and metadata from PPTX presentations."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []

    props = prs.core_properties
    for attr in ("author", "subject", "keywords", "comments", "title"):
        val = getattr(props, attr, None)
        if val:
            parts.append(f"{attr}: {val}")

    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"[slide {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def _parse_eml(data: bytes) -> str:
    """Extract text and headers from EML messages."""
    msg = BytesParser(policy=policy.default).parsebytes(data)
    parts = []

    for header in (
        "subject",
        "from",
        "to",
        "cc",
        "x-author",
        "x-subject-tag",
        "x-keywords",
        "x-description",
    ):
        value = msg.get(header)
        if value:
            parts.append(f"{header}: {value}")

    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            payload = part.get_payload(decode=True)
            if not payload:
                continue
            if ctype == "text/plain":
                parts.append(payload.decode(part.get_content_charset() or "utf-8", errors="replace"))
            elif ctype == "text/html":
                soup = BeautifulSoup(payload, "html.parser")
                parts.append(soup.get_text(separator="\n"))
                for img in soup.find_all("img"):
                    src = img.get("src")
                    if src:
                        parts.append(src)
    else:
        payload = msg.get_payload(decode=True) or b""
        parts.append(payload.decode(msg.get_content_charset() or "utf-8", errors="replace"))

    return "\n".join(parts)


def _parse_html(data: bytes) -> str:
    """Extract text from HTML, including meta tags."""
    soup = BeautifulSoup(data, "html.parser")
    parts = []
    # Extract meta tag content (intentionally vulnerable)
    for meta in soup.find_all("meta"):
        content = meta.get("content", "")
        if content:
            parts.append(content)
    parts.append(soup.get_text(separator="\n"))
    return "\n".join(parts)


def _ocr_image(data: bytes) -> str:
    """Run OCR on an image."""
    img = Image.open(io.BytesIO(data))
    return pytesseract.image_to_string(img)


def chunk_text(text: str) -> list[str]:
    """Split text into overlapping chunks."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        chunks.append(text[start:end])
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return [c.strip() for c in chunks if c.strip()]


def embed_and_store(chunks: list[str]):
    """Embed chunks and store in ChromaDB."""
    if not chunks:
        return
    embeddings = []
    for chunk in chunks:
        response = ollama_client.embed(model=OLLAMA_EMBED_MODEL, input=chunk)
        embeddings.append(response["embeddings"][0])
    ids = [f"chunk_{i}" for i in range(len(chunks))]
    collection.add(documents=chunks, embeddings=embeddings, ids=ids)


def query_documents(question: str, n_results: int = 3) -> list[str]:
    """Query ChromaDB for relevant chunks."""
    response = ollama_client.embed(model=OLLAMA_EMBED_MODEL, input=question)
    results = collection.query(
        query_embeddings=[response["embeddings"][0]],
        n_results=n_results,
    )
    return results["documents"][0] if results["documents"] else []


def ingest(file_bytes: bytes, filename: str) -> dict:
    """Full ingestion pipeline: parse -> chunk -> embed -> store."""
    global last_extracted_text, last_chunks

    text = parse_document(file_bytes, filename)
    last_extracted_text = text

    chunks = chunk_text(text)
    last_chunks = chunks

    embed_and_store(chunks)

    return {
        "filename": filename,
        "extracted_length": len(text),
        "num_chunks": len(chunks),
    }
