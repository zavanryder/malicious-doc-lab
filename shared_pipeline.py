"""Shared document processing pipeline used by both demo applications."""

import io
from email import policy
from email.parser import BytesParser
from pathlib import Path

import chromadb
import fitz  # PyMuPDF
import ollama
from bs4 import BeautifulSoup
from docx import Document
from PIL import Image
import pytesseract

MAX_XLSX_SCAN_ROWS = 20000
MAX_XLSX_SCAN_COLS = 256


class PipelineState:
    """Stateful parse -> chunk -> embed -> store pipeline."""

    def __init__(
        self,
        *,
        chroma_collection: str,
        ollama_base_url: str,
        ollama_embed_model: str,
        chunk_size: int,
        chunk_overlap: int,
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.ollama_embed_model = ollama_embed_model
        self.chroma_collection = chroma_collection
        self.chroma_client = chromadb.Client()
        self.collection = self.chroma_client.get_or_create_collection(
            name=chroma_collection
        )
        self.ollama_client = ollama.Client(host=ollama_base_url)
        self.last_extracted_text = ""
        self.last_chunks: list[str] = []

    def reset_state(self) -> None:
        """Clear all ingested data."""
        self.last_extracted_text = ""
        self.last_chunks = []
        self.chroma_client.delete_collection(name=self.chroma_collection)
        self.collection = self.chroma_client.get_or_create_collection(
            name=self.chroma_collection
        )

    def parse_document(self, file_bytes: bytes, filename: str) -> str:
        """Extract text from a document based on file extension."""
        suffix = Path(filename).suffix.lower()
        if suffix == ".pdf":
            return self._parse_pdf(file_bytes)
        if suffix == ".docx":
            return self._parse_docx(file_bytes)
        if suffix == ".xlsx":
            return self._parse_xlsx(file_bytes)
        if suffix == ".pptx":
            return self._parse_pptx(file_bytes)
        if suffix == ".eml":
            return self._parse_eml(file_bytes)
        if suffix in (".html", ".htm"):
            return self._parse_html(file_bytes)
        if suffix in (".csv", ".txt", ".md"):
            return file_bytes.decode("utf-8", errors="replace")
        if suffix in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
            return self._ocr_image(file_bytes)
        return file_bytes.decode("utf-8", errors="replace")

    def _parse_pdf(self, data: bytes) -> str:
        """Extract text from PDF, including OCR on embedded images."""
        doc = fitz.open(stream=data, filetype="pdf")
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
            for img_info in page.get_images():
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                if base_image:
                    img = Image.open(io.BytesIO(base_image["image"]))
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        text_parts.append(ocr_text)
        return "\n".join(text_parts)

    def _parse_docx(self, data: bytes) -> str:
        """Extract text from DOCX, including metadata."""
        doc = Document(io.BytesIO(data))
        parts = []
        props = doc.core_properties
        for attr in ("author", "subject", "keywords", "description", "title"):
            value = getattr(props, attr, None)
            if value:
                parts.append(f"{attr}: {value}")
        for para in doc.paragraphs:
            parts.append(para.text)
        return "\n".join(parts)

    def _parse_xlsx(self, data: bytes) -> str:
        """Extract text and metadata from XLSX workbooks."""
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []

        props = wb.properties
        for attr in ("creator", "subject", "keywords", "description", "title"):
            value = getattr(props, attr, None)
            if value:
                parts.append(f"{attr}: {value}")

        for sheet in wb.worksheets:
            parts.append(f"[sheet] {sheet.title}")
            max_row = min(sheet.max_row or 0, MAX_XLSX_SCAN_ROWS)
            max_col = min(sheet.max_column or 0, MAX_XLSX_SCAN_COLS)
            if max_row < 1 or max_col < 1:
                continue
            for row in sheet.iter_rows(
                min_row=1,
                max_row=max_row,
                min_col=1,
                max_col=max_col,
                values_only=True,
            ):
                values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if values:
                    parts.append(" | ".join(values))
        return "\n".join(parts)

    def _parse_pptx(self, data: bytes) -> str:
        """Extract text and metadata from PPTX presentations."""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts = []

        props = prs.core_properties
        for attr in ("author", "subject", "keywords", "comments", "title"):
            value = getattr(props, attr, None)
            if value:
                parts.append(f"{attr}: {value}")

        for index, slide in enumerate(prs.slides, start=1):
            parts.append(f"[slide {index}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)

    def _parse_eml(self, data: bytes) -> str:
        """Extract text and headers from EML messages."""
        msg = BytesParser(policy=policy.default).parsebytes(data)
        parts = []

        for header in (
            "subject",
            "from",
            "to",
            "cc",
            "x-author",
            "x-subject-tag",
            "x-keywords",
            "x-description",
        ):
            value = msg.get(header)
            if value:
                parts.append(f"{header}: {value}")

        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                payload = part.get_payload(decode=True)
                if not payload:
                    continue
                if ctype == "text/plain":
                    parts.append(
                        payload.decode(
                            part.get_content_charset() or "utf-8",
                            errors="replace",
                        )
                    )
                elif ctype == "text/html":
                    soup = BeautifulSoup(payload, "html.parser")
                    parts.append(soup.get_text(separator="\n"))
                    for img in soup.find_all("img"):
                        src = img.get("src")
                        if src:
                            parts.append(src)
        else:
            payload = msg.get_payload(decode=True) or b""
            parts.append(
                payload.decode(
                    msg.get_content_charset() or "utf-8",
                    errors="replace",
                )
            )

        return "\n".join(parts)

    def _parse_html(self, data: bytes) -> str:
        """Extract text from HTML, including meta tags."""
        soup = BeautifulSoup(data, "html.parser")
        parts = []
        for meta in soup.find_all("meta"):
            content = meta.get("content", "")
            if content:
                parts.append(content)
        parts.append(soup.get_text(separator="\n"))
        return "\n".join(parts)

    def _ocr_image(self, data: bytes) -> str:
        """Run OCR on an image."""
        img = Image.open(io.BytesIO(data))
        return pytesseract.image_to_string(img)

    def chunk_text(self, text: str) -> list[str]:
        """Split text into overlapping chunks."""
        chunks = []
        start = 0
        while start < len(text):
            end = start + self.chunk_size
            chunks.append(text[start:end])
            start += self.chunk_size - self.chunk_overlap
        return [chunk.strip() for chunk in chunks if chunk.strip()]

    def embed_and_store(self, chunks: list[str]) -> None:
        """Embed chunks and store them in ChromaDB."""
        if not chunks:
            return
        embeddings = []
        for chunk in chunks:
            response = self.ollama_client.embed(
                model=self.ollama_embed_model,
                input=chunk,
            )
            embeddings.append(response["embeddings"][0])
        ids = [f"chunk_{index}" for index in range(len(chunks))]
        self.collection.add(documents=chunks, embeddings=embeddings, ids=ids)

    def query_documents(self, question: str, n_results: int = 3) -> list[str]:
        """Query ChromaDB for relevant chunks."""
        response = self.ollama_client.embed(
            model=self.ollama_embed_model,
            input=question,
        )
        results = self.collection.query(
            query_embeddings=[response["embeddings"][0]],
            n_results=n_results,
        )
        return results["documents"][0] if results["documents"] else []

    def ingest(self, file_bytes: bytes, filename: str) -> dict:
        """Full ingestion pipeline: parse -> chunk -> embed -> store."""
        text = self.parse_document(file_bytes, filename)
        self.last_extracted_text = text

        chunks = self.chunk_text(text)
        self.last_chunks = chunks

        self.embed_and_store(chunks)

        return {
            "filename": filename,
            "extracted_length": len(text),
            "num_chunks": len(chunks),
        }
