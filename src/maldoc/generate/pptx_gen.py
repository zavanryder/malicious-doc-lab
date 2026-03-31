"""PPTX document generation with python-pptx."""

from pathlib import Path

from maldoc.attacks.base import AttackResult


def generate_pptx(attack_result: AttackResult, title: str, output_path: Path) -> Path:
    """Generate a PPTX presentation containing the attack payload."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title:
        slide.shapes.title.text = title

    body_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(1.2),
        Inches(9.0),
        Inches(4.5),
    )
    body_tf = body_box.text_frame
    body_tf.clear()
    body_tf.word_wrap = True
    body_tf.text = attack_result.visible_content

    if attack_result.metadata:
        props = prs.core_properties
        props.author = attack_result.metadata.get("author", "")
        props.subject = attack_result.metadata.get("subject", "")
        props.keywords = attack_result.metadata.get("keywords", "")
        props.comments = attack_result.metadata.get("description", "")

    hints = attack_result.format_hints or {}
    hidden_text = None
    if attack_result.technique == "white_on_white":
        hidden_text = attack_result.hidden_content
    elif attack_result.technique == "off_page":
        hidden_text = attack_result.hidden_content
    elif hints.get("obfuscated_variants"):
        hidden_text = hints["obfuscated_variants"][0]
    elif hints.get("obfuscated_payload"):
        hidden_text = hints["obfuscated_payload"]

    if not hidden_text:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    if attack_result.technique == "off_page":
        hidden_box = slide.shapes.add_textbox(
            Inches(12),
            Inches(8),
            Inches(2),
            Inches(1),
        )
    else:
        hidden_box = slide.shapes.add_textbox(
            Inches(0.6),
            Inches(6.8),
            Inches(8.8),
            Inches(0.5),
        )
    hidden_tf = hidden_box.text_frame
    hidden_tf.text = hidden_text
    for paragraph in hidden_tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(1 if attack_result.technique == "white_on_white" else 2)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
